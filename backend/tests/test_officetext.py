"""OOXML 로컬 고속 추출(officetext) 테스트.

라이브러리는 선택 의존성(extras: office) — 미설치면 실추출 테스트를 건너뛰고,
폴백(None) 계약은 설치 여부와 무관하게 검증한다.
"""

from __future__ import annotations

import pytest

from app import officetext


def test_broken_files_return_none(tmp_path):
    """깨진/DRM 래핑(비 zip) 파일은 예외 없이 None — COM 폴백 계약."""
    for name in ("a.docx", "b.xlsx", "c.pptx"):
        p = tmp_path / name
        p.write_bytes(b"DRM-wrapped, not a zip")
    assert officetext.extract_docx(str(tmp_path / "a.docx")) is None
    assert officetext.extract_xlsx(str(tmp_path / "b.xlsx")) is None
    assert officetext.extract_pptx(str(tmp_path / "c.pptx")) is None


def test_extract_docx_paragraphs_and_tables(tmp_path):
    docx = pytest.importorskip("docx")
    d = docx.Document()
    d.add_paragraph("HBM 시장 분석 개요")
    table = d.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "매출"
    table.rows[0].cells[1].text = "100"
    path = tmp_path / "r.docx"
    d.save(str(path))

    text = officetext.extract_docx(str(path))
    assert "HBM 시장 분석 개요" in text
    assert "매출\t100" in text


def test_extract_xlsx_matches_com_format(tmp_path):
    openpyxl = pytest.importorskip("openpyxl")
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "실적"
    ws["A1"], ws["B1"] = "매출", "QoQ"
    ws["A2"], ws["B2"] = 100, 3.2
    path = tmp_path / "f.xlsx"
    wb.save(str(path))

    text = officetext.extract_xlsx(str(path))
    assert "# 실적" in text          # COM 추출기와 같은 시트 헤더 형식
    assert "매출\tQoQ" in text
    assert "100\t3.2" in text


def test_extract_pptx_matches_com_format(tmp_path):
    pptx = pytest.importorskip("pptx")
    from pptx.util import Inches

    pres = pptx.Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "분기 하이라이트"
    path = tmp_path / "d.pptx"
    pres.save(str(path))

    text = officetext.extract_pptx(str(path))
    assert "--- slide 1 ---" in text  # COM 추출기와 같은 슬라이드 구분 형식
    assert "분기 하이라이트" in text


def test_textless_pptx_returns_none(tmp_path):
    """텍스트 없는 프레젠테이션은 None — 슬라이드 마커만으로 '본문 있음' 오판 방지."""
    pptx = pytest.importorskip("pptx")
    pres = pptx.Presentation()
    pres.slides.add_slide(pres.slide_layouts[6])
    path = tmp_path / "blank.pptx"
    pres.save(str(path))
    assert officetext.extract_pptx(str(path)) is None
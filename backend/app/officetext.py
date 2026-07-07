"""OOXML(docx/xlsx/pptx) 텍스트 추출 — Office 기동 없는 로컬 고속 경로.

OOXML 은 zip+XML 이라 python-docx / openpyxl / python-pptx 로 직접 파싱한다
(COM 의 초 단위 앱 기동·문서 열기 대비 밀리초 단위). 선택 의존성(extras: office).

폴백 계약: 미설치·파싱 실패·텍스트 없음이면 None 을 반환하고 호출부(COM 워커)가
Word/Excel/PowerPoint COM 으로 폴백한다. DRM 래핑 문서는 zip 시그니처가 깨져 있어
로컬 파서가 자연히 실패한다 — 즉 "일반 문서는 로컬, DRM 문서만 COM"이 자동 성립.

출력 형식은 COM 추출기와 맞춘다(시트 "# 이름"+탭 행, 슬라이드 "--- slide n ---") —
어느 경로로 추출되든 색인·검색 결과가 동일하도록.

이미지 처리: OCR(extras: ocr)이 활성이면 문서에 내장된 이미지(media/*)의 텍스트를
'[이미지 텍스트]' 블록으로 본문 뒤에 붙인다 — 스캔 캡처·표 이미지의 내용이
색인에서 유실되지 않도록. 개수·크기 상한으로 비용을 묶는다.
"""

from __future__ import annotations

from . import imagetext

_IMG_MIN_BYTES = 10 * 1024        # 아이콘·로고 등 노이즈 제외
_IMG_MAX_BYTES = 8 * 1024 * 1024  # 초대형 이미지 제외(비용 한도)
_IMG_MAX_COUNT = 10               # 문서당 OCR 이미지 상한
_IMG_EXTS = (".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff", ".gif")


def _embedded_image_text(path: str) -> str | None:
    """OOXML 내장 이미지(media/*)를 OCR/VLM 으로 텍스트화. 비활성/이미지 없음 → None."""
    if not imagetext.active():
        return None
    try:
        import zipfile  # noqa: PLC0415

        parts: list[str] = []
        with zipfile.ZipFile(path) as z:
            names = [n for n in z.namelist()
                     if "/media/" in n and n.lower().endswith(_IMG_EXTS)]
            for name in names[:_IMG_MAX_COUNT]:
                data = z.read(name)
                if not (_IMG_MIN_BYTES <= len(data) <= _IMG_MAX_BYTES):
                    continue
                text = imagetext.image_text(data)
                if text:
                    parts.append(text)
        return "\n".join(parts) or None
    except Exception:
        return None


def _with_image_text(path: str, body: str) -> str | None:
    """본문 뒤에 내장 이미지 OCR 텍스트를 붙인다. 둘 다 비면 None(COM 폴백)."""
    img = _embedded_image_text(path)
    if img:
        body = f"{body}\n[이미지 텍스트]\n{img}" if body else f"[이미지 텍스트]\n{img}"
    return body or None


def extract_docx(path: str) -> str | None:
    """docx 본문+표 텍스트. 미설치/실패/빈 문서 → None(COM 폴백)."""
    try:
        from docx import Document  # noqa: PLC0415
    except ImportError:
        return None
    try:
        doc = Document(path)
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = "\t".join(c.text.strip() for c in row.cells)
                if cells.strip():
                    parts.append(cells)
        return _with_image_text(path, "\n".join(parts))
    except Exception:
        return None


def extract_xlsx(path: str) -> str | None:
    """xlsx/xlsm 셀 텍스트(수식은 캐시된 값). 미설치/실패/빈 문서 → None(COM 폴백)."""
    try:
        from openpyxl import load_workbook  # noqa: PLC0415
    except ImportError:
        return None
    try:
        # read_only: 스트리밍 파싱(대용량 시트 메모리 안전), data_only: 수식 결과값
        wb = load_workbook(path, read_only=True, data_only=True)
        try:
            parts: list[str] = []
            for ws in wb.worksheets:
                rows = [
                    "\t".join("" if c is None else str(c) for c in row)
                    for row in ws.iter_rows(values_only=True)
                ]
                body = "\n".join(r for r in rows if r.strip())
                if body:
                    parts.append(f"# {ws.title}")
                    parts.append(body)
            return _with_image_text(path, "\n".join(parts))
        finally:
            wb.close()
    except Exception:
        return None


def extract_pptx(path: str) -> str | None:
    """pptx 슬라이드 텍스트. 미설치/실패/빈 문서 → None(COM 폴백)."""
    try:
        from pptx import Presentation  # noqa: PLC0415
    except ImportError:
        return None
    try:
        pres = Presentation(path)
        parts: list[str] = []
        any_text = False
        for idx, slide in enumerate(pres.slides, start=1):
            parts.append(f"--- slide {idx} ---")
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text)
                    any_text = True
        return _with_image_text(path, "\n".join(parts) if any_text else "")
    except Exception:
        return None
